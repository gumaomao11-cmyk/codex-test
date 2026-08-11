# -*- coding: utf-8 -*-
"""Optimize SMB Solution-Aipick1.pptx

Fixes:
  - TOC labels (slide 2): casing + alignment with actual sections
  - Slide 4: numbering spacing
  - Slide 6: double spaces, 'zones ,etc', inprocess, etc.
  - Slide 7: chapter title 'Product list and Topology' -> 'AiPick Solution'
  - Slide 8: 'Topology&' -> 'Topology &', 'Advance:' -> 'Advanced:'
  - Slides 10/11: 'Smartsearch' -> 'Smart Search'
  - Standardize 'Aipick' -> 'AiPick' (skipping model-name occurrences that must keep
    the original spelling if they map to a real product line).
  - Slide 14: 'Aipick cameras' -> 'AiPick cameras'
  - Append a new 'Why AiPick / Summary' slide.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy

SRC = r'D:\codex项目\shop\aipick_working.pptx'
DST = r'D:\codex项目\shop\aipick_optimized.pptx'


def set_textframe_text(tf, new_text):
    paragraphs = tf.paragraphs
    if not paragraphs:
        tf.text = new_text
        return
    first_p = paragraphs[0]
    runs = first_p.runs
    if runs:
        first_r = runs[0]
        first_r.text = new_text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first_p.text = new_text
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)


def replace_in_textframe(tf, replacements):
    for p in tf.paragraphs:
        runs = p.runs
        full = ''.join(r.text for r in runs) if runs else p.text
        new_full = full
        for old, new in replacements:
            if old in new_full:
                new_full = new_full.replace(old, new)
        if new_full != full:
            if runs:
                first = runs[0]
                first.text = new_full
                for r in runs[1:]:
                    r._r.getparent().remove(r._r)
            else:
                p.text = new_full




def merge_and_replace_paragraphs(tf, new_text):
    """For multi-paragraph text that should become a single line, merge all
    paragraphs into one and apply new_text on the first run."""
    paragraphs = tf.paragraphs
    if not paragraphs:
        tf.text = new_text
        return
    first_p = paragraphs[0]
    runs = first_p.runs
    if runs:
        first = runs[0]
        first.text = new_text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        first_p.text = new_text
    for p in paragraphs[1:]:
        p._p.getparent().remove(p._p)


def main():
    prs = Presentation(SRC)
    slides = list(prs.slides)
    print(f'Source slides: {len(slides)}')

    # === Slide 1: title slide brand name ===
    for sh in slides[0].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == 'Aipick':
            set_textframe_text(sh.text_frame, 'AiPick')
            print("  [S1 title] 'Aipick' -> 'AiPick'")


    # === Slide 2: TOC labels ===
    label_map = {
        'Solution benefits': 'Solution Benefits',
        'Product list and topology': 'AiPick Solution',
        'Successful Case': 'Successful Cases',
    }
    for sh in slides[1].shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t in label_map:
                set_textframe_text(sh.text_frame, label_map[t])
                print(f'  [TOC] {t!r} -> {label_map[t]!r}')

    # === Slide 4: numbering spacing ===
    for sh in slides[3].shapes:
        if sh.has_text_frame:
            replace_in_textframe(sh.text_frame, [
                ('1.Property loss', '1. Property loss'),
                ('3.Prolonged response time', '3. Prolonged response time'),
            ])

    # === Slide 6: Benefits — fix double space, punctuation, inprocess ===
    for sh in slides[5].shapes:
        if sh.has_text_frame:
            replace_in_textframe(sh.text_frame, [
                ('Benefits of  AI-powered Surveillance', 'Benefits of AI-powered Surveillance'),
                ('zones ,etc', 'zones, etc'),
                ('inprocess and', 'in-process and'),
            ])

    # === Slide 7: chapter title and AiPick casing ===
    for sh in slides[6].shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            if 'Product list' in t and 'Topology' in t:
                merge_and_replace_paragraphs(sh.text_frame, 'AiPick Solution')
            else:
                replace_in_textframe(sh.text_frame, [
                    ('Ai-Pick ', 'AiPick '),
                    ('Ai-Pick\n', 'AiPick\n'),
                    ('Ai-Pick,', 'AiPick,'),
                ])


    # === Slide 8: Topology title + Advance + spacing ===
    for sh in slides[7].shapes:
        if sh.has_text_frame:
            replace_in_textframe(sh.text_frame, [
                ('Aipick Solution Topology& Operational Flow',
                 'AiPick Solution Topology & Operational Flow'),
                ('Topology&', 'Topology &'),
                ('Advance:', 'Advanced:'),
                ('1.VERTICAL SYSTEM TOPOLOGY', '1. VERTICAL SYSTEM TOPOLOGY'),
                ('2.WORK FLOW', '2. WORK FLOW'),
                ('Aipick Dome', 'AiPick Dome'),
                ('Aipick Bullet', 'AiPick Bullet'),
                ('Aipick Eyeball', 'AiPick Eyeball'),
                ('Aipick Outdoor PTZ', 'AiPick Outdoor PTZ'),
                ('Aipick by NVR', 'AiPick by NVR'),
                ('Aipick by Camera', 'AiPick by Camera'),
            ])

    # === Slides 7-14: generic Aipick->AiPick pass for remaining body text ===
    generic = [('Aipick', 'AiPick')]
    for idx in (6, 7, 8, 9, 10, 11, 12, 13):
        for sh in slides[idx].shapes:
            if sh.has_text_frame:
                replace_in_textframe(sh.text_frame, generic)
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        replace_in_textframe(cell.text_frame, generic)


    # (catches everything we did not list explicitly above)
    global_pass = [
        ('Aipick cameras', 'AiPick cameras'),
        ('Aipick-supported', 'AiPick-supported'),
        ('Aipick by NVR', 'AiPick by NVR'),
        ('Aipick by Camera', 'AiPick by Camera'),
        ('Aipick Bullet', 'AiPick Bullet'),
        ('Aipick Dome', 'AiPick Dome'),
        ('Aipick Eyeball', 'AiPick Eyeball'),
        ('Aipick Outdoor PTZ', 'AiPick Outdoor PTZ'),
        ('Ai-Pick ', 'AiPick '),
        ('Ai-Pick\n', 'AiPick\n'),
        ('Ai-Pick,', 'AiPick,'),
        ('Ai-Pick.', 'AiPick.'),
        ('Smartsearch', 'Smart Search'),
    ]
    for idx in (8, 9, 10, 11, 12, 13):
        s = slides[idx]
        for sh in s.shapes:
            if sh.has_text_frame:
                replace_in_textframe(sh.text_frame, global_pass)
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        replace_in_textframe(cell.text_frame, global_pass)

    # === Append Summary slide ===
    template_slide = slides[12]  # 04 divider
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() in ('blank', '空白'):
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    new_slide = prs.slides.add_slide(blank_layout)
    for sh in template_slide.shapes:
        new_el = deepcopy(sh._element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # Rewrite the original "Successful Cases" heading to "Why AiPick"
    for sh in new_slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == 'Successful Cases':
            set_textframe_text(sh.text_frame, 'Why AiPick')

    # Remove the inherited "04" big number, then add a "05" badge in its place
    to_remove = []
    for sh in new_slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == '04':
            to_remove.append(sh)
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)

    badge = new_slide.shapes.add_textbox(
        Emu(int(1.16 * 914400)), Emu(int(0.0 * 914400)),
        Emu(int(5.5 * 914400)), Emu(int(5.0 * 914400)),
    )
    p = badge.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '05'
    run.font.name = 'Arimo Bold'
    run.font.size = Pt(242)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0F, 0x33, 0x66)

    sub = new_slide.shapes.add_textbox(
        Emu(int(9.5 * 914400)), Emu(int(3.4 * 914400)),
        Emu(int(10.0 * 914400)), Emu(int(0.8 * 914400)),
    )
    p = sub.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Key Takeaways'
    run.font.name = 'Arimo'
    run.font.size = Pt(22)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    bullets_tb = new_slide.shapes.add_textbox(
        Emu(int(1.4 * 914400)), Emu(int(5.0 * 914400)),
        Emu(int(17.0 * 914400)), Emu(int(5.5 * 914400)),
    )
    tf = bullets_tb.text_frame
    tf.word_wrap = True
    bullets = [
        ('24/7 complete monitoring',
         'Continuous coverage eliminates blind spots across corridors, lobbies and production zones.'),
        ('Real-time detection & alarm',
         'AI cameras flag intrusions, unauthorized access and high-risk events the moment they happen.'),
        ('Instant video retrieval',
         'AiPick one-click search across cameras, NVRs and the central platform - minutes, not hours.'),
        ('Spatiotemporal evidence chains',
         'Linkage across cameras builds a clear evidence trail for in-process and post-event review.'),
        ('Proven in real deployments',
         'Validated by garment factory and furniture mall roll-outs with measurable response gains.'),
    ]
    for i, (head, body) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run1 = p.add_run()
        run1.text = f'\u2022 {head}  '
        run1.font.name = 'Arimo Bold'
        run1.font.size = Pt(18)
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(0x0F, 0x33, 0x66)
        run2 = p.add_run()
        run2.text = body
        run2.font.name = 'Arimo'
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)

    prs.save(DST)
    print(f'Saved: {DST}')
    print(f'Final slide count: {len(prs.slides)}')


if __name__ == '__main__':
    main()
